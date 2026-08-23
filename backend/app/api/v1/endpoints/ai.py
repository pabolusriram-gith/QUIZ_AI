import io
import uuid
import json
import logging
import re
from typing import List, Optional, Dict, Any
from datetime import datetime, timezone

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, status, Response, Request
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.api.deps import get_current_user, check_role
from app.database.session import get_db
from app.models.user import User
from app.models.quiz import Question
from app.core.ai_providers import AIProviderFactory
from app.services.ai_retrieval import AIRetrievalService, RAGPromptBuilder
from app.services.ai_generator_service import PipelineCoordinator
from app.config.settings import settings
from app.core.rate_limit import check_ai_rate_limit
from app.services.ai_metrics_service import ai_metrics_service


router = APIRouter()
logger = logging.getLogger("app.api.ai")


def log_rejected_request(user_id: Optional[str], client_ip: str, endpoint: str, provider: str, request_type: str, retry_after: int, remaining_quota: int):
    log_data = {
        "event": "AI_RATE_LIMIT_REJECTED",
        "user_id": user_id,
        "client_ip": client_ip,
        "endpoint": endpoint,
        "provider": provider,
        "request_type": request_type,
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "retry_after": retry_after,
        "remaining_quota": remaining_quota
    }
    logger.warning(json.dumps(log_data))


class AIRateLimitGuard:
    def __init__(self, limit_type: str):
        self.limit_type = limit_type

    async def __call__(
        self,
        request: Request,
        question_count: int = Form(5),
        existing_questions: Optional[str] = Form(None),
        custom_prompt: Optional[str] = Form(None),
        current_user: User = Depends(get_current_user)
    ):
        # 1. Admin bypass
        if current_user and current_user.role == "admin":
            return

        # 2. Determine actual limit type
        actual_type = self.limit_type
        if self.limit_type == "generate":
            is_regeneration = (question_count == 1) and (existing_questions is not None or "replace this one" in (custom_prompt or ""))
            if is_regeneration:
                actual_type = "regenerate"

        # 3. Resolve rate-limiting key: authenticated user_id if present, client IP address if guest.
        user_id = str(current_user.id) if current_user else None
        client_ip = request.client.host if request.client else "unknown"
        key_id = user_id if user_id else client_ip

        # 4. Check rate limit
        is_allowed, retry_after = check_ai_rate_limit(key_id, actual_type)
        if not is_allowed:
            ai_metrics_service.increment_rejected()
            
            provider = "auto"
            try:
                form_data = await request.form()
                provider = form_data.get("provider", "auto")
            except Exception:
                pass

            log_rejected_request(
                user_id=user_id,
                client_ip=client_ip,
                endpoint=request.url.path,
                provider=provider,
                request_type=actual_type,
                retry_after=retry_after,
                remaining_quota=0
            )

            raise HTTPException(
                status_code=status.HTTP_429_TOO_MANY_REQUESTS,
                detail="Too many AI requests. Please wait a moment before trying again.",
                headers={"Retry-After": str(retry_after)}
            )

# Dynamic file parsing helpers
def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        import pypdf  # type: ignore[import-untyped]
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
        return text
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="The PDF parsing library 'pypdf' is not installed in the backend environment."
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to parse PDF document: {str(e)}"
        )

def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        import docx  # type: ignore[import-untyped]
        doc = docx.Document(io.BytesIO(file_bytes))
        text = ""
        for para in doc.paragraphs:
            if para.text:
                text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        return text
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="The Word parsing library 'python-docx' is not installed in the backend environment."
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to parse Word document: {str(e)}"
        )

def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        import pptx  # type: ignore[import-untyped]
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="The PowerPoint parsing library 'python-pptx' is not installed in the backend environment."
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to parse PowerPoint document: {str(e)}"
        )
def get_jaccard_similarity(s1: str, s2: str) -> float:
    w1 = set(s1.lower().split())
    w2 = set(s2.lower().split())
    if not w1 or not w2:
        return 0.0
    return len(w1.intersection(w2)) / len(w1.union(w2))


def validate_question(q: Dict[str, Any], validated_questions: List[Dict[str, Any]], existing_texts: List[str]) -> tuple[bool, str]:
    if not isinstance(q, dict):
        return False, "Not a dictionary"
    
    text = q.get("text", "").strip()
    if not text or len(text) < 15:
        return False, f"Question text too short or empty: '{text[:20]}...'"
        
    q_type = q.get("question_type", "").strip()
    if not q_type:
        return False, "Missing question_type"
        
    # Blacklisted Generic Questions check
    blacklist = [
        r"what\s+is\s+chatgpt",
        r"what\s+is\s+claude",
        r"what\s+is\s+gemini",
        r"expand\s+ai\b",
        r"what\s+is\s+machine\s+learning",
        r"who\s+developed\s+python",
        r"what\s+is\s+python\b",
        r"what\s+is\s+html\b",
        r"what\s+is\s+css\b",
        r"who\s+is\s+the\s+father\s+of"
    ]
    for pattern in blacklist:
        if re.search(pattern, text.lower()):
            return False, f"Matches generic question blacklist: '{pattern}'"
            
    # Duplicate checks within current batch
    for val_q in validated_questions:
        if get_jaccard_similarity(text, val_q.get("text", "")) > 0.5:
            return False, "Too similar to another question in the current batch"
            
    # Duplicate checks against database
    if existing_texts:
        for ext in existing_texts:
            if get_jaccard_similarity(text, ext) > 0.5:
                return False, "Too similar to an existing question in the Question Bank"
                
    # Options Validation
    options = q.get("options", [])
    if not isinstance(options, list):
        return False, "Options must be a list"
        
    if q_type == "multiple_choice":
        if len(options) != 4:
            return False, f"MCQ must have exactly 4 options, found {len(options)}"
        correct_count = sum(1 for o in options if o.get("is_correct"))
        if correct_count != 1:
            return False, f"MCQ must have exactly 1 correct option, found {correct_count}"
            
    elif q_type == "multiple_select":
        if len(options) != 4:
            return False, f"Multiple select must have exactly 4 options, found {len(options)}"
        correct_count = sum(1 for o in options if o.get("is_correct"))
        if correct_count < 1 or correct_count > 3:
            return False, f"Multiple select must have between 1 and 3 correct options, found {correct_count}"
            
    elif q_type == "true_false":
        if len(options) != 2:
            return False, f"True/False must have exactly 2 options, found {len(options)}"
        opt_texts = {o.get("text", "").lower().strip() for o in options}
        if not ("true" in opt_texts and "false" in opt_texts):
            return False, "True/False options must be 'True' and 'False'"
        correct_count = sum(1 for o in options if o.get("is_correct"))
        if correct_count != 1:
            return False, f"True/False must have exactly 1 correct option, found {correct_count}"
            
    elif q_type in ["fill_in_the_blank", "short_answer"]:
        if len(options) < 1:
            return False, "Fill in the blank/Short answer must have at least 1 correct value option"
        for o in options:
            if not o.get("text", "").strip():
                return False, "Fill in the blank/Short answer has empty option text"
            if not o.get("is_correct"):
                return False, "All option values in Fill in the blank/Short answer must be correct answers (is_correct=true)"
                
    # Check option duplicates within the question
    opt_set = set()
    for o in options:
        opt_txt = o.get("text", "").strip().lower()
        if opt_txt in opt_set:
            return False, f"Duplicate option text found within the question: '{opt_txt}'"
        opt_set.add(opt_txt)
        
    # Explanations check
    exp = q.get("explanation", "").strip()
    if not exp or len(exp) < 20:
        return False, "Explanation field is missing or too short (must be at least 20 chars)"
        
    return True, "OK"


def clean_metadata_text(text: str) -> str:
    if not text:
        return ""
    # Strip common prompt phrases
    prefixes = [
        "generate a quiz", "generate quiz", "generate questions", "generate a", "generate",
        "create a quiz", "create quiz", "create questions", "create",
        "based on the context of", "based on context", "based on the", "based on",
        "according to the", "according to",
        "using the context of", "using the", "using",
        "you are a", "you are",
        "quiz on", "questions on", "test on"
    ]
    cleaned = text.strip()
    while True:
        matched = False
        lower_cleaned = cleaned.lower()
        for p in prefixes:
            if lower_cleaned.startswith(p):
                cleaned = cleaned[len(p):].strip()
                cleaned = re.sub(r"^[\s:,\-._+]+", "", cleaned).strip()
                matched = True
                break
        if not matched:
            break
            
    if cleaned:
        cleaned = cleaned[0].upper() + cleaned[1:]
    return cleaned


@router.post("/generate", status_code=status.HTTP_200_OK, dependencies=[Depends(AIRateLimitGuard("generate"))])
async def generate_questions(
    request: Request,
    response: Response,
    question_count: int = Form(5),
    difficulty: str = Form("medium"),
    language: str = Form("en"),
    bloom_levels: str = Form(""), # Comma-separated list
    question_types: str = Form(""), # Comma-separated list
    topic: Optional[str] = Form(None),
    text: Optional[str] = Form(None),
    course_outcomes: Optional[str] = Form(None), # Comma-separated
    question_distribution: Optional[str] = Form(None),
    question_quality: str = Form("balanced"),
    quiz_style: str = Form("mixed"),
    custom_prompt: Optional[str] = Form(None),
    provider: str = Form("auto"),
    model_name: Optional[str] = Form(None),
    existing_questions: Optional[str] = Form(None),
    files: Optional[List[UploadFile]] = File(None),
    wrap_response: bool = Form(False),
    marks_mode: str = Form("default"),
    default_marks: int = Form(1),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> Any:
    """
    Generate questions from Topic, pasted text, or uploaded files.
    """
    # 1. Parse and index documents (RAG Indexing Phase)
    retrieval_service = AIRetrievalService()
    
    total_extracted_text = ""
    uploaded_filenames = []
    has_files_uploaded = False

    if files:
        has_files_uploaded = True
        
        ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt", ".txt", ".raw", ".csv", ".md"}
        ALLOWED_MIME_TYPES = {
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
            "text/plain",
            "text/csv",
            "text/markdown",
            "application/octet-stream"
        }
        MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
        
        for file in files:
            import os
            import re
            
            raw_filename = file.filename or "uploaded_document"
            safe_basename = os.path.basename(raw_filename)
            sanitized_name = re.sub(r"[^a-zA-Z0-9._-]", "_", safe_basename)
            
            ext = os.path.splitext(sanitized_name.lower())[1]
            if ext not in ALLOWED_EXTENSIONS:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Unsupported file format for file: '{raw_filename}'. Allowed formats: PDF, DOCX, PPT/PPTX, TXT, CSV, MD."
                )
            
            mime = file.content_type
            if mime and mime.lower() not in ALLOWED_MIME_TYPES:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid MIME type '{mime}' for file '{raw_filename}'."
                )

            try:
                file.file.seek(0, os.SEEK_END)
                file_size = file.file.tell()
                file.file.seek(0, os.SEEK_SET)
            except Exception:
                file_bytes_temp = await file.read()
                file_size = len(file_bytes_temp)
                file.file = io.BytesIO(file_bytes_temp)
            
            if file_size > MAX_FILE_SIZE:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"File '{raw_filename}' exceeds the maximum allowed size of 10MB."
                )
                
            file_bytes = await file.read()
            
            # File signature magic bytes validation
            is_valid_sig = False
            if ext == ".pdf":
                is_valid_sig = file_bytes.startswith(b"%PDF")
            elif ext in {".docx", ".pptx"}:
                is_valid_sig = file_bytes.startswith(b"PK\x03\x04")
            elif ext == ".ppt":
                is_valid_sig = file_bytes.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1")
            elif ext in {".txt", ".raw", ".csv", ".md"}:
                try:
                    file_bytes.decode("utf-8")
                    is_valid_sig = True
                except UnicodeDecodeError:
                    is_valid_sig = False
                    
            if not is_valid_sig:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"File signature validation failed for '{raw_filename}'."
                )
                
            content = ""
            if ext == ".pdf":
                content = extract_text_from_pdf(file_bytes)
            elif ext == ".docx":
                content = extract_text_from_docx(file_bytes)
            elif ext in {".ppt", ".pptx"}:
                content = extract_text_from_pptx(file_bytes)
            elif ext in {".txt", ".raw", ".csv", ".md"}:
                content = file_bytes.decode("utf-8")
                
            clean_content = content.strip()
            if clean_content:
                total_extracted_text += clean_content + "\n"
                unique_doc_name = f"{uuid.uuid4()}_{sanitized_name}"
                uploaded_filenames.append(unique_doc_name)

    if text and text.strip():
        pasted_clean = text.strip()
        total_extracted_text += pasted_clean + "\n"
        unique_doc_name = f"{uuid.uuid4()}_pasted_text"
        uploaded_filenames.append(unique_doc_name)

    # Enforce parsed text limit to prevent memory exhaustion / prompt injection attacks
    MAX_TEXT_LIMIT = 500000
    if len(total_extracted_text) > MAX_TEXT_LIMIT:
        total_extracted_text = total_extracted_text[:MAX_TEXT_LIMIT]

    if has_files_uploaded:
        non_whitespace_count = sum(1 for c in total_extracted_text if not c.isspace())
        if non_whitespace_count < 20:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Document text extraction failed. The uploaded file(s) contain little or no readable text."
            )

    request_id = str(uuid.uuid4())

    # 2. Determine execution mode and index documents (RAG Indexing Phase)
    CONTEXT_THRESHOLD = 50000
    has_document_context = bool(total_extracted_text.strip())
    document_names = None
    final_context = ""

    if has_document_context:
        text_length = len(total_extracted_text)
        if text_length <= CONTEXT_THRESHOLD:
            # DIRECT CONTEXT MODE
            if settings.DEBUG or settings.ENVIRONMENT == "development":
                logger.warning(f"Direct Context Mode active (Extracted length: {text_length} <= {CONTEXT_THRESHOLD})")
            final_context = f"[EXTRACTED DOCUMENT CONTENT]\n{total_extracted_text}\n"
            if topic:
                final_context = f"Topic/Context Prompt: {topic}\n\n" + final_context
        else:
            # RAG MODE
            if settings.DEBUG or settings.ENVIRONMENT == "development":
                logger.warning(f"RAG Mode active (Extracted length: {text_length} > {CONTEXT_THRESHOLD})")
            
            rag_doc_name = f"{uuid.uuid4()}_combined_document"
            await retrieval_service.index_document(db, rag_doc_name, total_extracted_text)
            document_names = [rag_doc_name]
            
            # Retrieve chunks for single-question regeneration if needed
            search_query = topic or custom_prompt or "assessment questions"
            retrieved_chunks = await retrieval_service.retrieve_relevant_chunks(
                db, search_query, limit=5, document_names=document_names
            )
            base_prompt = ""
            if topic:
                base_prompt += f"Topic/Context Prompt: {topic}\n"
            final_context = RAGPromptBuilder.build_prompt(base_prompt, retrieved_chunks)
    else:
        # TOPIC-BASED MODE
        if topic:
            final_context = f"Topic/Context Prompt: {topic}\n"

    if not has_document_context and not topic:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Please provide a Topic, Pasted Text, or upload at least one valid Document file."
        )

    # 3. Parse input arguments and delegate to PipelineCoordinator
    q_types_list = [t.strip() for t in question_types.split(",") if t.strip()]
    if not q_types_list:
        q_types_list = ["multiple_choice"]
        
    bloom_list = [b.strip() for b in bloom_levels.split(",") if b.strip()]
    if not bloom_list:
        bloom_list = ["Understand", "Apply"]

    co_list = [c.strip() for c in course_outcomes.split(",") if c.strip()] if course_outcomes else []

    user_id = str(current_user.id) if current_user else None
    is_regeneration = (question_count == 1) and (existing_questions is not None or "replace this one" in (custom_prompt or ""))

    if is_regeneration:
        # Premium Regeneration Pipeline
        exist_list = []
        if existing_questions:
            try:
                exist_list = json.loads(existing_questions)
                if not isinstance(exist_list, list):
                    exist_list = []
            except Exception as e:
                logger.warning(f"Failed to parse existing_questions JSON: {e}")
                
        # Parse the original question text from the custom prompt if possible
        original_q_text = ""
        original_match = re.search(r'replace this one:\s*"(.*?)"', custom_prompt or "", re.DOTALL)
        if original_match:
            original_q_text = original_match.group(1).strip()
            
        # Extract the clean, non-generic original quiz topic/subject
        original_quiz_topic = ""
        for candidate in [topic, course_outcomes, original_q_text]:
            if candidate:
                cand_clean = candidate.strip().lower()
                if cand_clean not in ["", "general", "general knowledge", "project alpha", "general topic", "general concept", "assessment questions"]:
                    original_quiz_topic = candidate.strip()
                    break

        target_topic = original_quiz_topic or topic or "General"

        original_question = {
            "text": original_q_text or "Original placeholder question text",
            "difficulty": difficulty,
            "bloom_level": bloom_list[0] if bloom_list else "Understand",
            "question_type": q_types_list[0] if q_types_list else "multiple_choice",
            "topic": target_topic,
            "learning_objective": target_topic,
            "marks_mode": marks_mode,
            "marks": default_marks
        }
        
        result = await PipelineCoordinator.regenerate_question(
            db=db,
            original_question=original_question,
            existing_questions=exist_list,
            provider=provider,
            model_name=model_name,
            quiz_style=quiz_style,
            question_quality=question_quality,
            final_context=final_context,
            request_id=request_id,
            user_id=user_id
        )
        if user_id:
            ai_metrics_service.track_user_usage(user_id, "regenerate", provider)
        return result
    else:
        # Standard quiz generation with targeted failures replacement loop
        questions = await PipelineCoordinator.generate_quiz(
            db=db,
            question_count=question_count,
            difficulty=difficulty,
            language=language,
            bloom_levels=bloom_list,
            question_types=q_types_list,
            topic=topic,
            course_outcomes=co_list,
            question_distribution=question_distribution,
            question_quality=question_quality,
            quiz_style=quiz_style,
            custom_prompt=custom_prompt,
            provider=provider,
            model_name=model_name,
            request_id=request_id,
            document_names=document_names,
            total_extracted_text=total_extracted_text,
            marks_mode=marks_mode,
            default_marks=default_marks,
            user_id=user_id
        )
        if has_document_context and len(questions) < question_count:
            warning_msg = f"Insufficient document content. Generated {len(questions)} of {question_count} requested questions."
            response.headers["X-Warning"] = warning_msg
            logger.warning(warning_msg)
        user_id = str(current_user.id) if current_user else None
        if user_id:
            ai_metrics_service.track_user_usage(user_id, "generate", provider)
            
        if wrap_response:
            cleaned_topic = clean_metadata_text(topic or "General Assessment")
            metadata = {
                "title": f"{cleaned_topic} Quiz",
                "subject": cleaned_topic,
                "department": "Academics",
                "description": f"AI-generated assessment covering {cleaned_topic}."
            }
            return {
                "metadata": metadata,
                "questions": questions
            }
        return questions


@router.post("/check-duplicate")
async def check_duplicates(
    payload: Dict[str, Any],
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Compare generated questions with the existing Question Bank (questions table) 
    using token-based Jaccard similarity.
    """
    questions_in = payload.get("questions", [])
    if not questions_in:
        return {"warnings": []}

    # Fetch existing questions
    result = await db.execute(select(Question.text))
    existing_texts = [r[0] for r in result.all()]

    warnings = []

    def get_jaccard_similarity(s1: str, s2: str) -> float:
        w1 = set(s1.lower().split())
        w2 = set(s2.lower().split())
        if not w1 or not w2:
            return 0.0
        return len(w1.intersection(w2)) / len(w1.union(w2))

    for idx, q in enumerate(questions_in):
        q_text = q.get("text", "")
        if not q_text:
            continue
            
        max_sim = 0.0
        match_text = ""
        
        for ext in existing_texts:
            sim = get_jaccard_similarity(q_text, ext)
            if sim > max_sim:
                max_sim = sim
                match_text = ext
                
        # Similarity threshold (> 70% match)
        if max_sim > 0.7:
            warnings.append({
                "index": idx,
                "similarity": round(max_sim, 2),
                "existing_text": match_text,
                "warning": f"Similar question already exists in Question Bank ({int(max_sim * 100)}% match)."
            })

    return {"warnings": warnings}


@router.post("/regenerate-question", status_code=status.HTTP_200_OK)
async def regenerate_single_question(
    request: Request,
    # The original question's core constraints
    original_question_json: str = Form(..., description="JSON-encoded original question dict with at minimum: text, difficulty, question_type, bloom_level, topic, marks"),
    # Other questions in the quiz for deduplication
    existing_questions_json: Optional[str] = Form(None, description="JSON-encoded list of other question texts or dicts to avoid duplication"),
    # Quiz context for RAG and AI generation
    quiz_subject: Optional[str] = Form(None, description="Quiz subject (e.g. 'Computer Science')"),
    quiz_topic: Optional[str] = Form(None, description="Quiz generation topic/prompt used during initial generation"),
    # AI provider config
    provider: str = Form("auto"),
    model_name: Optional[str] = Form(None),
    question_quality: str = Form("balanced"),
    quiz_style: str = Form("mixed"),
    # Marks config – must match quiz's marks settings
    marks_mode: str = Form("default"),
    default_marks: int = Form(1),
    # Novelty & override params – used when difficulty/type changes trigger regeneration
    uniqueness_hint: Optional[str] = Form(None, description="Millisecond timestamp or random token to guarantee a novel question on every call"),
    target_difficulty_override: Optional[str] = Form(None, description="New difficulty to use instead of the original question's difficulty"),
    target_type_override: Optional[str] = Form(None, description="New question_type to use instead of the original question's type"),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> Any:
    """
    Dedicated single-question regeneration endpoint using the full QuizVerse RAG pipeline.

    Regenerates exactly ONE question to replace an existing AI-generated question.
    - Uses PipelineCoordinator.regenerate_question (full pipeline with QuestionProcessor)
    - Performs RAG retrieval on quiz_topic/quiz_subject for relevant context
    - Enforces the original question's difficulty, type, and bloom_level
    - Respects marks_mode (auto: easy=1, medium=2, hard=5 / default: uses default_marks)
    - Does NOT regenerate any other questions
    - Rate-limited under "regenerate" quota
    """
    # ── 1. Rate-limit check ────────────────────────────────────────────────────
    if current_user.role != "admin":
        user_id_str = str(current_user.id)
        client_ip = request.client.host if request.client else "unknown"
        key_id = user_id_str
        is_allowed, retry_after = check_ai_rate_limit(key_id, "regenerate")
        if not is_allowed:
            ai_metrics_service.increment_rejected()
            log_rejected_request(
                user_id=user_id_str,
                client_ip=client_ip,
                endpoint=request.url.path,
                provider=provider,
                request_type="regenerate",
                retry_after=retry_after,
                remaining_quota=0
            )
            raise HTTPException(
                status_code=status.HTTP_429_TOO_MANY_REQUESTS,
                detail="Too many AI regeneration requests. Please wait a moment before trying again.",
                headers={"Retry-After": str(retry_after)}
            )

    # ── 2. Parse original question ─────────────────────────────────────────────
    try:
        original_question = json.loads(original_question_json)
        if not isinstance(original_question, dict):
            raise ValueError("original_question_json must be a JSON object")
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid original_question_json: {e}"
        )

    # ── 3. Parse existing questions list ───────────────────────────────────────
    existing_questions: List[Any] = []
    if existing_questions_json:
        try:
            parsed = json.loads(existing_questions_json)
            if isinstance(parsed, list):
                existing_questions = parsed
        except Exception as e:
            logger.warning(f"Failed to parse existing_questions_json: {e}. Continuing without deduplication list.")

    # ── 4. Build RAG context for this question ─────────────────────────────────
    # Extract the clean, non-generic original quiz topic/subject
    original_quiz_topic = ""
    for candidate in [quiz_topic, quiz_subject, original_question.get("topic")]:
        if candidate:
            cand_clean = candidate.strip().lower()
            if cand_clean not in ["", "general", "general knowledge", "project alpha", "general topic", "general concept", "assessment questions"]:
                original_quiz_topic = candidate.strip()
                break

    topic = original_quiz_topic or original_question.get("topic") or quiz_topic or quiz_subject or "General"
    search_query = topic

    final_context = ""
    retrieval_service = AIRetrievalService()

    try:
        # Attempt RAG retrieval using embedded search query
        retrieved_chunks = await retrieval_service.retrieve_relevant_chunks(
            db, search_query, limit=5
        )
        if retrieved_chunks:
            base_prompt = ""
            if quiz_topic:
                base_prompt += f"Topic/Context Prompt: {quiz_topic}\n"
            elif quiz_subject:
                base_prompt += f"Subject: {quiz_subject}\n"
            final_context = RAGPromptBuilder.build_prompt(base_prompt, retrieved_chunks)
            logger.info(f"RAG retrieval successful for regeneration: {len(retrieved_chunks)} chunks retrieved")
        else:
            # No chunks found – fall back to topic-only context (safe fallback)
            if quiz_topic:
                final_context = f"Topic/Context Prompt: {quiz_topic}\n"
            elif quiz_subject:
                final_context = f"Subject: {quiz_subject}\n"
            logger.info("RAG retrieval returned 0 chunks. Using topic-only context as safe fallback.")
    except Exception as rag_err:
        logger.warning(f"RAG retrieval failed during regeneration: {rag_err}. Using topic-only context.")
        if quiz_topic:
            final_context = f"Topic/Context Prompt: {quiz_topic}\n"
        elif quiz_subject:
            final_context = f"Subject: {quiz_subject}\n"

    # ── 5. Enrich original_question with marks_mode for QuestionProcessor ─────
    # PipelineCoordinator.regenerate_question reads original_question["marks_mode"]
    # and original_question["marks"] to pass to QuestionProcessor.process_question.
    original_question["marks_mode"] = marks_mode
    if marks_mode == "auto":
        diff = (original_question.get("difficulty") or "medium").lower()
        if diff == "easy":
            original_question["marks"] = 1
        elif diff == "hard":
            original_question["marks"] = 5
        else:
            original_question["marks"] = 2
    else:
        # Use the quiz's default_marks unless the question already has a marks value
        if not original_question.get("marks"):
            original_question["marks"] = default_marks

    # Ensure topic is set for subject profile matching
    if not original_question.get("topic"):
        original_question["topic"] = topic

    # ── 5b. Apply target overrides (difficulty/type changes trigger full AI regen) ──
    # When the user changes difficulty or question_type in the editor, the frontend
    # sends the new value as target_difficulty_override / target_type_override.
    # We overwrite the relevant fields in original_question so PipelineCoordinator
    # generates with the NEW constraints rather than the stale original ones.
    if target_difficulty_override and target_difficulty_override.lower() in ["easy", "medium", "hard"]:
        original_question["difficulty"] = target_difficulty_override.lower()
        logger.info(f"[regenerate-question] difficulty override applied: {target_difficulty_override}")

    if target_type_override and target_type_override in [
        "multiple_choice", "multiple_select", "true_false", "fill_in_the_blank", "short_answer"
    ]:
        original_question["question_type"] = target_type_override
        logger.info(f"[regenerate-question] question_type override applied: {target_type_override}")

    # ── 5c. Inject uniqueness hint for novelty guarantee ──────────────────────
    if uniqueness_hint:
        original_question["uniqueness_hint"] = uniqueness_hint

    # ── 6. Generate request_id for tracing ────────────────────────────────────
    request_id = str(uuid.uuid4())

    logger.info(f"[regenerate-question] request_id={request_id} "
                f"difficulty={original_question.get('difficulty')} "
                f"type={original_question.get('question_type')} "
                f"bloom={original_question.get('bloom_level')} "
                f"marks_mode={marks_mode} "
                f"provider={provider} "
                f"uniqueness_hint={uniqueness_hint!r} "
                f"diff_override={target_difficulty_override!r} "
                f"type_override={target_type_override!r}")

    user_id = str(current_user.id) if current_user else None

    # ── 7. Call the existing PipelineCoordinator.regenerate_question ──────────
    result = await PipelineCoordinator.regenerate_question(
        db=db,
        original_question=original_question,
        existing_questions=existing_questions,
        provider=provider,
        model_name=model_name,
        quiz_style=quiz_style,
        question_quality=question_quality,
        final_context=final_context,
        request_id=request_id,
        user_id=user_id
    )

    # ── 8. Track usage metrics ────────────────────────────────────────────────
    if user_id:
        ai_metrics_service.track_user_usage(user_id, "regenerate", provider)

    # Return the single regenerated question
    return result[0] if result else {}


@router.post("/enhance", status_code=status.HTTP_200_OK, dependencies=[Depends(AIRateLimitGuard("enhance"))])
async def enhance_prompt(
    request: Request,
    prompt: str = Form(...),
    provider: str = Form("auto"),
    model_name: Optional[str] = Form(None),
    current_user: User = Depends(get_current_user)
) -> Any:
    """
    Enhance a simple user prompt into a structured, professional quiz-generation instruction.
    """
    from app.config.settings import settings
    from app.core.ai_providers import resolve_provider, is_gemini_configured, is_groq_configured, is_openai_configured, AIProviderConfigurationError
    import google.generativeai as genai
    import asyncio
    
    system_instruction = (
        "You are an expert AI prompt engineer. Rewrite the user's quiz generation topic/prompt "
        "into a highly detailed, professional, and structured instruction for a quiz generator. "
        "Make sure to specify a good balance of subjects, difficulty, bloom's taxonomy categories, and formatting. "
        "Do NOT output any conversational text, explanations, or quotes. Output ONLY the rewritten prompt string itself."
    )
    
    try:
        resolved_name = resolve_provider(provider)
    except AIProviderConfigurationError as config_err:
        if provider != "auto":
            raise HTTPException(status_code=400, detail=str(config_err))
        else:
            raise HTTPException(
                status_code=503,
                detail="All AI providers are currently unavailable. Please try again later."
            )
            
    # Priority list for "auto" mode, or explicit provider
    providers_to_try = []
    if resolved_name == "auto":
        for prov in settings.AUTO_PROVIDER_ORDER:
            if prov == "gemini" and is_gemini_configured():
                providers_to_try.append("gemini")
            elif prov == "groq" and is_groq_configured():
                providers_to_try.append("groq")
            elif prov == "openai" and is_openai_configured():
                providers_to_try.append("openai")
            elif prov == "mock" and settings.ENABLE_MOCK_PROVIDER:
                providers_to_try.append("mock")
        if settings.ENABLE_MOCK_PROVIDER and "mock" not in providers_to_try:
            providers_to_try.append("mock")
    else:
        providers_to_try = [resolved_name]
        
    import time
    errors = []
    for p in providers_to_try:
        if p != "mock":
            from app.services.ai_health_service import ai_health_service
            state = ai_health_service.check_and_update_circuit(p)
            if state == "OPEN":
                logger.warning(f"Auto mode fallback (enhance): Skipping provider {p} because its circuit breaker is OPEN.")
                errors.append(f"{p} (circuit open): skipped")
                continue

        if p == "mock":
            start_t = time.time()
            words = [w.strip() for w in prompt.split(" ") if w.strip()]
            topic_name = " ".join(words[-3:]) if len(words) >= 3 else (prompt if prompt else "general subjects")
            mock_enhanced = (
                f"Generate 20 medium-difficulty conceptual and scenario-based questions covering {topic_name}, "
                "focusing on core methodologies, key definitions, real-world applications, and evaluation metrics. "
                "Ensure the questions test various levels of Bloom's Taxonomy (Understand, Apply, Analyze) "
                "and include clear, detailed explanations for the correct options."
            )
            gen_time = int((time.time() - start_t) * 1000)
            ai_metrics_service.increment_success("mock")
            ai_metrics_service.record_latency("mock", gen_time)
            
            user_id = str(current_user.id) if current_user else None
            if user_id:
                ai_metrics_service.track_user_usage(user_id, "enhance", "mock")
            return {"enhanced_prompt": mock_enhanced}
            
        elif p == "gemini":
            from app.core.ai_providers import execute_resilient_raw_call
            async def run_gemini():
                genai.configure(api_key=settings.GEMINI_API_KEY)
                model = genai.GenerativeModel(
                    model_name=model_name or "gemini-3.5-flash",
                    system_instruction=system_instruction
                )
                response = await asyncio.to_thread(model.generate_content, prompt)
                return response.text.strip()
            try:
                start_t = time.time()
                enhanced_text = await execute_resilient_raw_call("gemini", run_gemini)
                gen_time = int((time.time() - start_t) * 1000)
                
                ai_metrics_service.increment_success("gemini")
                ai_metrics_service.record_latency("gemini", gen_time)
                
                user_id = str(current_user.id) if current_user else None
                if user_id:
                    ai_metrics_service.track_user_usage(user_id, "enhance", "gemini")
                return {"enhanced_prompt": enhanced_text}
            except Exception as e:
                ai_metrics_service.increment_failure("gemini")
                logger.error(f"Gemini prompt enhancement failed: {e}")
                errors.append(f"Gemini: {e}")
                
        elif p == "groq":
            from app.core.ai_providers import execute_resilient_raw_call
            from openai import AsyncOpenAI
            async def run_groq():
                client = AsyncOpenAI(api_key=settings.GROQ_API_KEY, base_url="https://api.groq.com/openai/v1")
                response = await client.chat.completions.create(
                    model=model_name or "llama-3.3-70b-versatile",
                    messages=[
                        {"role": "system", "content": system_instruction},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7
                )
                return response.choices[0].message.content.strip()
            try:
                start_t = time.time()
                enhanced_text = await execute_resilient_raw_call("groq", run_groq)
                gen_time = int((time.time() - start_t) * 1000)
                
                ai_metrics_service.increment_success("groq")
                ai_metrics_service.record_latency("groq", gen_time)
                
                user_id = str(current_user.id) if current_user else None
                if user_id:
                    ai_metrics_service.track_user_usage(user_id, "enhance", "groq")
                return {"enhanced_prompt": enhanced_text}
            except Exception as e:
                ai_metrics_service.increment_failure("groq")
                logger.error(f"Groq prompt enhancement failed: {e}")
                errors.append(f"Groq: {e}")
                
        elif p == "openai":
            from app.core.ai_providers import execute_resilient_raw_call
            from openai import AsyncOpenAI
            async def run_openai():
                client = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
                response = await client.chat.completions.create(
                    model=model_name or "gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": system_instruction},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7
                )
                return response.choices[0].message.content.strip()
            try:
                start_t = time.time()
                enhanced_text = await execute_resilient_raw_call("openai", run_openai)
                gen_time = int((time.time() - start_t) * 1000)
                
                ai_metrics_service.increment_success("openai")
                ai_metrics_service.record_latency("openai", gen_time)
                
                user_id = str(current_user.id) if current_user else None
                if user_id:
                    ai_metrics_service.track_user_usage(user_id, "enhance", "openai")
                return {"enhanced_prompt": enhanced_text}
            except Exception as e:
                ai_metrics_service.increment_failure("openai")
                logger.error(f"OpenAI prompt enhancement failed: {e}")
                errors.append(f"OpenAI: {e}")
                
    # If all failed
    if provider == "auto":
        raise HTTPException(
            status_code=503,
            detail="All AI providers are currently unavailable. Please try again later."
        )
    else:
        is_unavailable = any("circuit open" in err.lower() or "unavailable" in err.lower() for err in errors)
        raise HTTPException(
            status_code=503 if is_unavailable else 400,
            detail=f"Prompt enhancement failed. Errors: {'; '.join(errors)}"
        )



@router.get("/providers", status_code=status.HTTP_200_OK)
async def get_enabled_providers(
    current_user: User = Depends(get_current_user)
) -> Any:
    """
    Get the list of active/enabled AI providers with metadata.
    """
    from app.config.settings import settings
    from app.core.ai_providers import is_gemini_configured, is_groq_configured, is_openai_configured

    providers = [
        {
            "id": "auto",
            "name": "Auto",
            "recommended": True
        }
    ]
    if is_gemini_configured():
        providers.append({
            "id": "gemini",
            "name": "Gemini",
            "recommended": False
        })
    if is_groq_configured():
        providers.append({
            "id": "groq",
            "name": "Groq",
            "recommended": False
        })
    if is_openai_configured():
        providers.append({
            "id": "openai",
            "name": "OpenAI",
            "recommended": False
        })
    if settings.ENABLE_MOCK_PROVIDER:
        providers.append({
            "id": "mock",
            "name": "Mock Engine",
            "recommended": False
        })
        
    return {"providers": providers}



@router.get("/metrics", status_code=status.HTTP_200_OK)
async def get_ai_metrics(
    include_traces: bool = False,
    trace_limit: int = 50,
    current_user: User = Depends(get_current_user)
) -> Any:
    """
    Retrieve current AI latency metrics and trace history snapshot.
    """
    return ai_metrics_service.get_internal_snapshot(
        include_traces=include_traces,
        trace_limit=trace_limit,
        user_id=str(current_user.id) if current_user.role != "admin" else None
    )


